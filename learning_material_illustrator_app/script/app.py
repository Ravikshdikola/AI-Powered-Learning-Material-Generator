import streamlit as st
import os
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
import base64

# --- CONFIG ---base64
st.set_page_config(page_title="AI Learning Material", layout="centered")
genai.configure(api_key="AIzaSyAAIIizm7e2TIPrKvbqVbryB24Yie3bnhI")

image_model = genai.GenerativeModel('gemini-2.5-flash-image-preview')
text_model = genai.GenerativeModel('gemini-2.5-flash')

# --- UI: Header ---
st.title("📚 AI-Powered Learning Material Generator")
st.markdown("Generate labeled diagrams and explanations for middle school topics using Gemini AI.")

# --- Step 0: Topic Selection + Manual Input ---
default_topics = ["Water Cycle", "Solar System", "Digestive System"]

# Multiselect with **no default selected**
topics_selected = st.multiselect("Select topics:", default_topics)

# Manual topics input as multi-line text area (one topic per line)
custom_topics_input = st.text_area("Add custom topics (one per line):")

# Process manual topics: split lines, strip, remove empty
custom_topics = [t.strip() for t in custom_topics_input.split('\n') if t.strip()]

# Combine and remove duplicates while keeping order
all_topics = list(dict.fromkeys(topics_selected + custom_topics))

if st.button("🔄 Generate Educational Content") and all_topics:
    st.info("Generating images and explanations. This may take a minute...")

    image_files = []

    progress = st.progress(0)
    total = len(all_topics)

    for i, topic in enumerate(all_topics, start=1):
        # --- Step 1: Generate image prompt ---
        img_prompt_request = f"""
        Generate an image prompt for an educational diagram of '{topic}'.
        1. The style should be simple, colorful, and clear.
        2. It must look like a labeled diagram, not abstract art or a story.
        3. Make sure labels are easy to read for middle school students.
        4. Avoid extra creativity or storytelling, focus on clarity and accuracy.
        Output only the image prompt.
        """
        img_prompt = text_model.generate_content(img_prompt_request).text.strip()

        # --- Step 2: Generate image ---
        response = image_model.generate_content(img_prompt)

        img_data = None
        for part in response.candidates[0].content.parts:
            if part.inline_data and part.inline_data.mime_type.startswith("image/"):
                img_data = part.inline_data.data
                break

        if img_data:
            filename = f"{topic.replace(' ', '_')}.png"
            with open(filename, "wb") as f:
                f.write(img_data)
            image_files.append((filename, topic, topic))
            st.image(filename, caption=f"{topic} Diagram", width=500)
        else:
            st.warning(f"⚠️ Image not generated for: {topic}")

        progress.progress(i / total)

    st.success("✅ All images generated!")

    # --- Step 3: Generate PowerPoint ---
    prs = Presentation()

    for filename, title, topic in image_files:
        # Slide 1: Image + Title (using slide title placeholder)
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])
        title_shape = slide1.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(24)
        slide1.shapes.add_picture(filename, Inches(1), Inches(1.5), height=Inches(3.5))

        # Slide 2: Explanation + Title (using slide title placeholder)
        explanation_prompt = f"""
        Write a simple, clear explanation for middle school students about: {topic}.
        1. Start with a one-sentence definition of {topic}.
        2. Then explain the key ideas in 3–4 short sentences.
        3. Keep it under 100 words.
        4. Use simple words, avoid storytelling tone.
        """
        explanation = text_model.generate_content(explanation_prompt).text.strip()

        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        title_shape2 = slide2.shapes.title
        title_shape2.text = f"{title} - Explanation"
        title_shape2.text_frame.paragraphs[0].font.size = Pt(24)

        text_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8), Inches(5))
        tf = text_box.text_frame
        tf.word_wrap = True
        para = tf.add_paragraph()
        para.text = explanation
        para.font.size = Pt(16)

    pptx_filename = "AI_Learning_Material.pptx"
    prs.save(pptx_filename)

    # --- Download Link ---
    def get_binary_file_downloader_html(file_path, file_label='File'):
        with open(file_path, 'rb') as f:
            data = f.read()
        b64 = base64.b64encode(data).decode()
        href = f'<a href="data:application/octet-stream;base64,{b64}" download="{file_path}">📥 Download {file_label}</a>'
        return href

    st.markdown("---")
    st.markdown(get_binary_file_downloader_html(pptx_filename, 'PowerPoint'), unsafe_allow_html=True)
    st.balloons()

else:
    st.info("Select at least one topic or add custom topics (one per line), then click 'Generate Educational Content'.")
